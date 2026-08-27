"""Extract searchable text from supported module document formats."""

from pathlib import Path
import re

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}
CHUNK_SIZE = 2000


def _clean_text(value: object) -> str:
    return re.sub(r"\s+", " ", str(value or "")).strip()


def _extract_pdf(path: Path) -> list[str]:
    reader = PdfReader(str(path))
    return [_clean_text(page.extract_text()) for page in reader.pages]


def _extract_docx(path: Path) -> list[str]:
    document = Document(str(path))
    fragments = [_clean_text(paragraph.text) for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            fragments.append(_clean_text(" | ".join(cell.text for cell in row.cells)))
    return fragments


def _extract_xlsx(path: Path) -> list[str]:
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    try:
        fragments = []
        for worksheet in workbook.worksheets:
            for row in worksheet.iter_rows(values_only=True):
                values = [_clean_text(value) for value in row]
                values = [value for value in values if value]
                if values:
                    fragments.append(" | ".join(values))
        return fragments
    finally:
        workbook.close()


def _extract_pptx(path: Path) -> list[str]:
    presentation = Presentation(str(path))
    fragments = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = _clean_text(shape.text)
                if text:
                    fragments.append(text)
    return fragments


def extract_module_text(path: str | Path) -> dict[str, object]:
    """Return normalized chunks and metadata for a supported document."""
    document_path = Path(path)
    extension = document_path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError("Unsupported module format. Use PDF, DOCX, XLSX, or PPTX.")

    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".pptx": _extract_pptx,
    }
    fragments = [_clean_text(fragment) for fragment in extractors[extension](document_path)]
    text = "\n".join(fragment for fragment in fragments if fragment)
    chunks = [text[index:index + CHUNK_SIZE] for index in range(0, len(text), CHUNK_SIZE)]

    return {
        "chunks": chunks,
        "content_extracted": bool(text),
        "warning": None if text else "No extractable text was found in this file.",
    }